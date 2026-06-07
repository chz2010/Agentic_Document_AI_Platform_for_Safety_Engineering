"""Create two black-and-white PowerPoint architecture decks.

The decks are intentionally minimal: black background, white text, white
outlines, large spacing, and short copy so all text stays inside boxes.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_DIR = Path(__file__).resolve().parents[1]
OUTPUT_DIR = PROJECT_DIR / "presentations"

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)

SLIDE_W = 13.333
SLIDE_H = 7.5


def new_deck() -> Presentation:
    deck = Presentation()
    deck.slide_width = Inches(SLIDE_W)
    deck.slide_height = Inches(SLIDE_H)
    return deck


def add_slide(deck: Presentation, title: str, kicker: str = ""):
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BLACK
    if kicker:
        add_text(slide, kicker.upper(), 0.55, 0.32, 11.8, 0.25, 8, bold=True)
    add_text(slide, title, 0.55, 0.72, 12.0, 0.75, 25, bold=True)
    add_line(slide, 0.55, 1.55, 12.75, 1.55, 1.0)
    return slide


def add_text(slide, value: str, x: float, y: float, w: float, h: float, size: int, bold: bool = False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    p = frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = value
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = WHITE
    return shape


def add_box(slide, title: str, lines: list[str], x: float, y: float, w: float, h: float, title_size: int = 13, body_size: int = 10):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLACK
    shape.line.color.rgb = WHITE
    shape.line.width = Pt(1.2)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = Inches(0.16)
    frame.margin_right = Inches(0.16)
    frame.margin_top = Inches(0.11)
    frame.margin_bottom = Inches(0.1)

    p = frame.paragraphs[0]
    p.text = title
    p.font.name = "Aptos"
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(7)

    for line in lines:
        p = frame.add_paragraph()
        p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(body_size)
        p.font.color.rgb = WHITE
        p.level = 0
        p.space_after = Pt(4)
    return shape


def add_metric_box(slide, label: str, value: str, x: float, y: float, w: float, h: float):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLACK
    shape.line.color.rgb = WHITE
    shape.line.width = Pt(1.2)
    add_text(slide, label.upper(), x + 0.16, y + 0.16, w - 0.32, 0.22, 8, bold=True)
    add_text(slide, value, x + 0.16, y + 0.52, w - 0.32, 0.38, 18, bold=True)


def add_line(slide, x1: float, y1: float, x2: float, y2: float, width: float = 1.2):
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    shape.line.color.rgb = WHITE
    shape.line.width = Pt(width)
    return shape


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float):
    shape = add_line(slide, x1, y1, x2, y2, 1.4)
    shape.line.end_arrowhead = True
    return shape


def footer(slide, label: str, page: int):
    add_text(slide, label, 0.55, 7.12, 9.8, 0.2, 8)
    add_text(slide, f"{page}", 12.35, 7.12, 0.4, 0.2, 8, align=PP_ALIGN.RIGHT)


def title_slide(deck: Presentation, title: str, subtitle: str, label: str):
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BLACK
    add_text(slide, label.upper(), 0.7, 0.6, 10.0, 0.28, 9, bold=True)
    add_text(slide, title, 0.7, 1.45, 11.7, 1.35, 31, bold=True)
    add_text(slide, subtitle, 0.75, 3.25, 9.8, 0.72, 15)
    add_box(slide, "Core message", ["Safety AI is valuable when it becomes an auditable engineering workflow.", "This platform shows that full path."], 0.75, 4.65, 5.55, 1.2, 14, 11)
    add_box(slide, "Portfolio signal", ["Backend architecture", "Document AI", "Requirements engineering", "Agent operations"], 6.75, 4.65, 5.55, 1.2, 14, 11)
    return slide


def build_project2_architecture_deck() -> Presentation:
    deck = new_deck()
    title_slide(
        deck,
        "Project 2 Architecture",
        "Agentic Document AI Platform for Safety Engineering",
        "Architecture narrative",
    )

    s = add_slide(deck, "Executive Summary", "Project 2")
    add_box(s, "What it is", ["A production-style backend platform.", "Turns safety documents into structured engineering outputs.", "Works as a standalone second project."], 0.75, 1.95, 3.8, 1.65, 14, 11)
    add_box(s, "What it proves", ["FastAPI backend design.", "Document AI and project RAG.", "Requirements extraction and quality scoring.", "Traceability and Agent Ops."], 4.85, 1.95, 3.8, 1.65, 14, 11)
    add_box(s, "Why it matters", ["The system is more than chat.", "It stores evidence, actions, evaluations, and reports.", "That is production thinking."], 8.95, 1.95, 3.6, 1.65, 14, 11)
    add_metric_box(s, "Operating model", "Upload -> Analyze -> Trace -> Report", 1.1, 4.8, 11.1, 0.95)
    footer(s, "Project 2 architecture", 2)

    s = add_slide(deck, "System Architecture", "Runtime view")
    add_box(s, "Streamlit UI", ["Workspace", "Upload", "Ask", "Reports"], 0.7, 2.35, 2.25, 1.25)
    add_box(s, "FastAPI", ["Typed endpoints", "Orchestration", "Pydantic outputs"], 3.75, 2.2, 2.35, 1.55)
    add_box(s, "PostgreSQL", ["Projects", "Runs", "Workflow", "Requirements"], 6.9, 1.65, 2.25, 1.25)
    add_box(s, "Chroma", ["Project chunks", "Embeddings", "Metadata filters"], 6.9, 3.35, 2.25, 1.25)
    add_box(s, "Model Layer", ["OpenAI", "Local model", "No-LLM mode"], 10.0, 2.2, 2.45, 1.55)
    add_arrow(s, 2.95, 2.95, 3.75, 2.95)
    add_arrow(s, 6.1, 2.65, 6.9, 2.15)
    add_arrow(s, 6.1, 3.1, 6.9, 3.95)
    add_arrow(s, 9.15, 3.0, 10.0, 3.0)
    footer(s, "Project 2 architecture", 3)

    s = add_slide(deck, "Core Data Model", "Persistence view")
    add_box(s, "Workspace layer", ["Project", "Document", "DocumentChunk"], 0.75, 2.0, 3.3, 1.35, 15, 12)
    add_box(s, "Engineering layer", ["Requirement", "TestCase", "TraceabilityLink"], 4.95, 2.0, 3.3, 1.35, 15, 12)
    add_box(s, "Operations layer", ["EvaluationRun", "AgentRunLog", "WorkflowItem"], 9.15, 2.0, 3.3, 1.35, 15, 12)
    add_box(s, "Design principle", ["Every output is tied back to project evidence.", "Every run is logged.", "Every gap can become workflow action."], 1.4, 4.45, 10.5, 1.25, 15, 12)
    footer(s, "Project 2 architecture", 4)

    s = add_slide(deck, "Document AI Pipeline", "From file to evidence")
    labels = [
        ("1", "Upload", "PDF, TXT, MD, CSV, DOCX"),
        ("2", "Parse", "Extract text and metadata"),
        ("3", "Chunk", "Project-specific evidence units"),
        ("4", "Embed", "OpenAI or local hash fallback"),
        ("5", "Retrieve", "Filtered by project ID"),
    ]
    x = 0.7
    for n, title, body in labels:
        add_metric_box(s, n, title, x, 2.2, 1.75, 0.9)
        add_box(s, "", [body], x, 3.35, 1.75, 0.85, 1, 11)
        if n != "5":
            add_arrow(s, x + 1.75, 2.65, x + 2.1, 2.65)
        x += 2.45
    add_box(s, "Recruiter signal", ["The system handles ingestion as a pipeline, not as a one-off prompt."], 2.0, 5.35, 9.25, 0.75, 14, 12)
    footer(s, "Project 2 architecture", 5)

    s = add_slide(deck, "Requirements Engineering Module", "Structured safety outputs")
    add_box(s, "Extract", ["Find shall-style requirements.", "Classify requirement type.", "Preserve evidence source."], 0.9, 2.0, 3.15, 1.55, 15, 11)
    add_box(s, "Evaluate", ["Score clarity.", "Score testability.", "Flag ambiguity and traceability gaps."], 5.05, 2.0, 3.15, 1.55, 15, 11)
    add_box(s, "Generate", ["Candidate requirements.", "Test cases.", "Traceability matrix."], 9.2, 2.0, 3.15, 1.55, 15, 11)
    add_box(s, "Output contract", ["Hazard -> Safety Goal -> Requirement -> Test Case -> Evidence -> Status"], 1.3, 4.75, 10.8, 0.95, 15, 13)
    footer(s, "Project 2 architecture", 6)

    s = add_slide(deck, "Agent Operations", "Production behavior")
    add_box(s, "Run logging", ["Agent run ID", "Model used", "Prompt version", "Tools used"], 0.75, 1.9, 2.75, 1.75, 14, 11)
    add_box(s, "Reliability", ["Confidence score", "Failure reason", "Human escalation", "Approval gate"], 3.95, 1.9, 2.75, 1.75, 14, 11)
    add_box(s, "Cost control", ["Prompt tokens", "Output tokens", "Latency", "Estimated cost"], 7.15, 1.9, 2.75, 1.75, 14, 11)
    add_box(s, "Dashboard", ["Success rate", "Escalation rate", "Quality trend", "Model comparison"], 10.35, 1.9, 2.2, 1.75, 14, 11)
    add_box(s, "Why this is strong", ["The platform treats AI as an operated service, not a black box."], 1.4, 5.0, 10.4, 0.85, 15, 12)
    footer(s, "Project 2 architecture", 7)

    s = add_slide(deck, "Deployment And Demo Story", "What to say")
    add_box(s, "Local demo", ["FastAPI backend", "Streamlit UI", "SQLite default", "Chroma local vector DB"], 0.8, 1.95, 3.55, 1.65, 15, 11)
    add_box(s, "Production path", ["Docker Compose", "PostgreSQL", "Vector DB service", "Optional Streamlit frontend"], 4.9, 1.95, 3.55, 1.65, 15, 11)
    add_box(s, "Domain adaptability", ["Automotive today.", "Railway-ready when licensed standards are provided.", "Educational context clearly labeled."], 9.0, 1.95, 3.55, 1.65, 15, 11)
    add_box(s, "Closing line", ["This project demonstrates backend engineering, safety-domain AI, structured outputs, and operational monitoring in one coherent platform."], 1.1, 4.95, 11.0, 0.95, 15, 12)
    footer(s, "Project 2 architecture", 8)
    return deck


def build_interaction_deck() -> Presentation:
    deck = new_deck()
    title_slide(
        deck,
        "Project 1 + Project 2 Interaction",
        "How the Safety Analyst and the Document AI Platform become one safety engineering workflow.",
        "System interaction narrative",
    )

    s = add_slide(deck, "The Portfolio Logic", "Two systems, one story")
    add_box(s, "Project 1", ["Autonomous Driving Safety Analyst", "Domain RAG assistant", "Explains standards and videos"], 0.95, 2.0, 4.55, 1.65, 15, 11)
    add_box(s, "Project 2", ["Agentic Document AI Platform", "Engineering workflow backend", "Turns documents into actions"], 7.8, 2.0, 4.55, 1.65, 15, 11)
    add_arrow(s, 5.65, 2.8, 7.65, 2.8)
    add_box(s, "Combined value", ["Knowledge retrieval becomes requirements review, traceability, workflow, and reports."], 1.7, 4.75, 9.9, 0.9, 15, 12)
    footer(s, "Project interaction", 2)

    s = add_slide(deck, "Boundary Between Projects", "No confusion")
    add_box(s, "Project 1 answers", ["What does this concept mean?", "What evidence exists in transcripts?", "How do standards ideas relate?"], 0.8, 2.0, 3.7, 1.75, 15, 11)
    add_box(s, "Project 2 answers", ["Is this requirement good?", "What traceability is missing?", "What test case is needed?"], 4.95, 2.0, 3.7, 1.75, 15, 11)
    add_box(s, "Human role", ["Curates notes.", "Approves context.", "Reviews compliance claims."], 9.1, 2.0, 3.25, 1.75, 15, 11)
    add_box(s, "Safe claim", ["Project 1 is knowledge support. Project 2 is engineering control."], 1.45, 5.0, 10.4, 0.85, 15, 12)
    footer(s, "Project interaction", 3)

    s = add_slide(deck, "Current Practical Workflow", "Works today")
    add_box(s, "1. Ingest", ["Project 1 ingests railway lecture transcripts."], 0.8, 2.2, 2.25, 1.15, 14, 11)
    add_box(s, "2. Ask", ["Project 1 explains RAMS concepts."], 3.55, 2.2, 2.25, 1.15, 14, 11)
    add_box(s, "3. Curate", ["Engineer saves approved notes."], 6.3, 2.2, 2.25, 1.15, 14, 11)
    add_box(s, "4. Review", ["Project 2 checks requirements."], 9.05, 2.2, 2.25, 1.15, 14, 11)
    add_arrow(s, 3.05, 2.78, 3.55, 2.78)
    add_arrow(s, 5.8, 2.78, 6.3, 2.78)
    add_arrow(s, 8.55, 2.78, 9.05, 2.78)
    add_box(s, "Important framing", ["Lecture transcripts are educational context, not official IEC/EN standard text."], 1.2, 5.0, 10.9, 0.85, 15, 12)
    footer(s, "Project interaction", 4)

    s = add_slide(deck, "Future FastAPI Integration", "Target design")
    add_box(s, "Project 2", ["Needs domain context during review."], 0.85, 2.2, 3.0, 1.15, 15, 11)
    add_box(s, "Project 1 API", ["Search transcripts and standards context."], 5.15, 1.7, 3.0, 1.15, 15, 11)
    add_box(s, "Evidence response", ["Summary", "Citations", "Limitations"], 5.15, 3.65, 3.0, 1.25, 15, 11)
    add_box(s, "Project 2 output", ["Requirement gaps", "Traceability", "Report"], 9.55, 2.55, 2.75, 1.35, 15, 11)
    add_arrow(s, 3.85, 2.75, 5.15, 2.25)
    add_arrow(s, 8.15, 2.25, 9.55, 3.05)
    add_arrow(s, 8.15, 4.2, 9.55, 3.45)
    footer(s, "Project interaction", 5)

    s = add_slide(deck, "Railway Use Case", "Codewerk-relevant angle")
    add_box(s, "Railway context", ["11 IEC 62278 / EN 50126 lecture transcripts.", "Used as public educational RAMS context."], 0.8, 1.95, 3.75, 1.55, 15, 11)
    add_box(s, "Requirements dataset", ["ETCS requirements document.", "Train control and supervision requirements."], 4.85, 1.95, 3.75, 1.55, 15, 11)
    add_box(s, "Review objective", ["Find safety-critical requirements.", "Identify missing verification evidence.", "Build traceability gaps."], 8.9, 1.95, 3.45, 1.55, 15, 11)
    add_box(s, "Result", ["A railway-flavored demo without claiming official railway compliance."], 1.25, 4.85, 10.7, 0.85, 15, 12)
    footer(s, "Project interaction", 6)

    s = add_slide(deck, "Recruiter Walkthrough", "How to present it")
    add_box(s, "Step 1", ["Show Project 1 answering RAMS concept questions from transcripts."], 0.8, 1.9, 2.8, 1.25, 14, 11)
    add_box(s, "Step 2", ["Show Project 2 ingesting ETCS requirements."], 3.95, 1.9, 2.8, 1.25, 14, 11)
    add_box(s, "Step 3", ["Extract and score requirements."], 7.1, 1.9, 2.35, 1.25, 14, 11)
    add_box(s, "Step 4", ["Show Agent Ops and report export."], 9.8, 1.9, 2.55, 1.25, 14, 11)
    add_box(s, "Message", ["You built both the AI knowledge layer and the production workflow layer."], 1.5, 4.85, 10.3, 0.95, 16, 12)
    footer(s, "Project interaction", 7)

    s = add_slide(deck, "Closing Architecture Claim", "What they should remember")
    add_box(s, "Single sentence", ["Project 1 retrieves safety knowledge; Project 2 turns that knowledge into auditable engineering decisions."], 1.1, 2.0, 11.1, 1.05, 18, 14)
    add_box(s, "Why it is strong", ["It shows RAG, backend APIs, document AI, requirements engineering, traceability, workflow, reporting, and AI operations."], 1.1, 3.75, 11.1, 1.2, 18, 14)
    footer(s, "Project interaction", 8)
    return deck


def main() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    architecture_path = OUTPUT_DIR / "Project_2_Architecture_Agentic_Document_AI_Platform.pptx"
    interaction_path = OUTPUT_DIR / "Project_1_Project_2_Interaction_Architecture.pptx"
    build_project2_architecture_deck().save(architecture_path)
    build_interaction_deck().save(interaction_path)
    print(architecture_path)
    print(interaction_path)


if __name__ == "__main__":
    main()
